import re
import tempfile
import unittest
from pathlib import Path


SAMPLE_MARKDOWN = """# Informe Ejecutivo

## Resumen
Texto con **énfasis** y una [fuente](https://example.com/reporte) verificable.

| Herramienta | Categoría | Prioridad |
| :--- | :--- | :--- |
| Ollama | Modelos locales | Alta |
| Chroma | RAG | Media |

- Privacidad local
- Costo operativo bajo

## Fuentes consultadas
1. https://example.com/reporte
"""


class DesignMdDeliverableTests(unittest.TestCase):
    def test_lint_design_md_detects_core_quality_issues(self):
        from k_zero_core.services.design_md import lint_design_md

        bad_design = """---
name: Bad
colors:
  surface: "#ffffff"
components:
  card:
    backgroundColor: "{colors.surface}"
    textColor: "{colors.missing}"
---
# Overview
# Overview
# Colors
"""

        findings = lint_design_md(bad_design)
        codes = {finding.code for finding in findings}

        self.assertIn("missing-primary", codes)
        self.assertIn("missing-typography", codes)
        self.assertIn("broken-ref", codes)
        self.assertIn("duplicate-section", codes)

    def test_markdown_cleanup_removes_raw_markdown_markers(self):
        from k_zero_core.services.design_md import limpiar_markdown_entregable

        cleaned = limpiar_markdown_entregable(SAMPLE_MARKDOWN)

        self.assertNotIn("**", cleaned)
        self.assertNotIn("###", cleaned)
        self.assertNotRegex(cleaned, re.compile(r"^\|", re.MULTILINE))
        self.assertIn("fuente (https://example.com/reporte)", cleaned)

    def test_docx_uses_native_headings_tables_and_bullets(self):
        from docx import Document
        from k_zero_core.core.tools import documents

        with tempfile.TemporaryDirectory() as tmp:
            original_exports = documents.EXPORTS_DIR
            documents.EXPORTS_DIR = Path(tmp)
            try:
                result = documents.crear_docx(SAMPLE_MARKDOWN, "informe", design_md_path="K-Zero ejecutivo")
                output = Path(result.split("Archivo creado: ", 1)[1].splitlines()[0])
                doc = Document(str(output))
            finally:
                documents.EXPORTS_DIR = original_exports
        paragraph_text = "\n".join(paragraph.text for paragraph in doc.paragraphs)

        self.assertEqual(doc.core_properties.title, "Informe Ejecutivo")
        self.assertEqual(doc.core_properties.author, "K-Zero")
        self.assertIn("K-Zero", doc.core_properties.keywords)
        self.assertIn("Entregable ejecutivo", paragraph_text)
        self.assertIn("Fuentes verificadas", paragraph_text)
        self.assertGreaterEqual(len(doc.tables), 1)
        self.assertIn("Informe Ejecutivo", paragraph_text)
        self.assertIn("Privacidad local", paragraph_text)
        self.assertNotIn("**", paragraph_text)
        self.assertNotRegex(paragraph_text, re.compile(r"^\|", re.MULTILINE))

    def test_pdf_uses_wrapped_platypus_content_without_raw_markdown(self):
        from pypdf import PdfReader
        from k_zero_core.core.tools import documents

        long_markdown = SAMPLE_MARKDOWN + "\n" + ("Este párrafo largo debe fluir sin cortarse en el borde de la página. " * 18)
        with tempfile.TemporaryDirectory() as tmp:
            original_exports = documents.EXPORTS_DIR
            documents.EXPORTS_DIR = Path(tmp)
            try:
                result = documents.crear_pdf(long_markdown, "resumen", design_md_path="")
                output = Path(result.split("Archivo creado: ", 1)[1].splitlines()[0])
                text = "\n".join(page.extract_text() or "" for page in PdfReader(str(output)).pages)
            finally:
                documents.EXPORTS_DIR = original_exports

        self.assertIn("Informe Ejecutivo", text)
        self.assertIn("Ollama", text)
        self.assertIn("Este párrafo largo debe fluir", text)
        self.assertNotIn("**", text)
        self.assertNotIn("| Herramienta |", text)

    def test_xlsx_styles_headers_freezes_and_filters(self):
        from openpyxl import load_workbook
        from k_zero_core.core.tools import documents

        with tempfile.TemporaryDirectory() as tmp:
            original_exports = documents.EXPORTS_DIR
            documents.EXPORTS_DIR = Path(tmp)
            try:
                result = documents.crear_xlsx(SAMPLE_MARKDOWN, "matriz", design_md_path="")
                output = Path(result.split("Archivo creado: ", 1)[1].splitlines()[0])
                wb = load_workbook(str(output))
            finally:
                documents.EXPORTS_DIR = original_exports

        ws = wb.active

        self.assertEqual(ws.freeze_panes, "A2")
        self.assertIsNotNone(ws.auto_filter.ref)
        self.assertEqual(ws["A1"].value, "Herramienta")
        self.assertEqual(ws["A1"].font.bold, True)
        self.assertNotEqual(ws["A1"].fill.fgColor.rgb, "00000000")
        wb.close()

    def test_pptx_creates_multiple_native_slides_without_raw_markdown(self):
        from pptx import Presentation
        from k_zero_core.core.tools import documents

        with tempfile.TemporaryDirectory() as tmp:
            original_exports = documents.EXPORTS_DIR
            documents.EXPORTS_DIR = Path(tmp)
            try:
                result = documents.crear_pptx(SAMPLE_MARKDOWN, "presentacion", design_md_path="")
                output = Path(result.split("Archivo creado: ", 1)[1].splitlines()[0])
                prs = Presentation(str(output))
            finally:
                documents.EXPORTS_DIR = original_exports

        all_text = "\n".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text")
        )

        self.assertGreaterEqual(len(prs.slides), 3)
        self.assertIn("Informe Ejecutivo", all_text)
        self.assertIn("Resumen", all_text)
        self.assertNotIn("**", all_text)
        self.assertNotIn("| Herramienta |", all_text)

    def test_design_tools_are_registered(self):
        from k_zero_core.core.tools import get_all_tools

        names = [tool.__name__ for tool in get_all_tools()]

        self.assertIn("validar_design_md", names)
        self.assertIn("crear_design_md", names)
        self.assertIn("previsualizar_estilo_entregable", names)
        self.assertIn("limpiar_markdown_entregable", names)


if __name__ == "__main__":
    unittest.main()
