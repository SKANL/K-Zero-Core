"""Tools documentales para lectura, creación y edición segura por copia."""
from __future__ import annotations

import csv
import json
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from k_zero_core.core.config import DATA_DIR
from k_zero_core.core.tool_safety import resolve_safe_path
from k_zero_core.services.document_reader import extract_text
from k_zero_core.services.design_md import (
    aplicar_diseno_entregable,
    clean_inline_text,
    parse_markdown_blocks,
)


EXPORTS_DIR = DATA_DIR / "exports"


def _export_path(name: str, suffix: str) -> Path:
    EXPORTS_DIR.mkdir(parents=True, exist_ok=True)
    clean = "".join(char for char in name if char.isalnum() or char in ("-", "_")).strip() or "entregable"
    if not suffix.startswith("."):
        suffix = "." + suffix
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return EXPORTS_DIR / f"{clean}_{timestamp}{suffix}"


def _markdown_lines(markdown: str) -> list[str]:
    return [line.rstrip() for line in markdown.splitlines()]


def _read_text_like(path: Path, max_chars: int) -> str:
    text = path.read_text(encoding="utf-8", errors="replace")
    return text[:max_chars] + ("\n[...truncado...]" if len(text) > max_chars else "")


def _created_result(path: Path, output_type: str, suggested_name: str) -> str:
    return (
        f"Archivo creado: {path}\n"
        f"Tipo: {output_type}\n"
        f"Nombre sugerido: {suggested_name}\n"
        "Estado: creado"
    )


def _theme_color(deliverable: dict[str, Any], token_path: str, fallback: str) -> str:
    design = deliverable["design"]
    value = design.token(token_path, fallback)
    return value if isinstance(value, str) and value.startswith("#") else fallback


def _hex_rgb(hex_color: str) -> tuple[int, int, int]:
    clean = hex_color.strip().lstrip("#")
    return tuple(int(clean[index : index + 2], 16) for index in (0, 2, 4))


def _first_heading(blocks: list[dict[str, Any]], default: str) -> str:
    for block in blocks:
        if block["type"] == "heading" and block.get("level") == 1:
            return block["text"]
    return default


def _table_rows(block: dict[str, Any]) -> list[list[str]]:
    return [block["headers"], *block["rows"]]


def _set_docx_cell_shading(cell: Any, fill: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    tc_pr = cell._tc.get_or_add_tcPr()
    shading = OxmlElement("w:shd")
    shading.set(qn("w:fill"), fill.strip("#"))
    tc_pr.append(shading)


def _set_docx_font_color(run: Any, hex_color: str) -> None:
    from docx.shared import RGBColor

    red, green, blue = _hex_rgb(hex_color)
    run.font.color.rgb = RGBColor(red, green, blue)


def _configure_docx_styles(doc: Any, primary: str) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt

    normal = doc.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.12

    for style_name, size, bold in (("Title", 24, False), ("Subtitle", 11, False), ("Heading 1", 14, True), ("Heading 2", 12, True)):
        style = doc.styles[style_name]
        style.font.name = "Aptos"
        style.font.size = Pt(size)
        style.font.bold = bold
        style.paragraph_format.space_before = Pt(10 if "Heading" in style_name else 0)
        style.paragraph_format.space_after = Pt(5)
        if style_name in {"Title", "Subtitle"}:
            style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT


def _add_docx_cover(doc: Any, title: str, subtitle: str, primary: str) -> None:
    from datetime import date
    from docx.enum.text import WD_BREAK
    from docx.shared import Pt

    title_paragraph = doc.add_paragraph(style="Title")
    title_run = title_paragraph.add_run(title)
    title_run.font.size = Pt(24)
    _set_docx_font_color(title_run, primary)

    if subtitle:
        subtitle_paragraph = doc.add_paragraph(style="Subtitle")
        subtitle_paragraph.add_run(subtitle)

    meta = doc.add_paragraph()
    meta.add_run("Entregable ejecutivo").bold = True
    meta.add_run(f"\nGenerado por K-Zero\n{date.today().isoformat()}")

    separator = doc.add_paragraph()
    separator_run = separator.add_run("─" * 58)
    _set_docx_font_color(separator_run, primary)

    separator.add_run().add_break(WD_BREAK.PAGE)


def _add_docx_sources(doc: Any, sources: list[str], primary: str) -> None:
    heading = doc.add_heading("Fuentes verificadas" if sources else "Fuentes no verificadas", level=1)
    for run in heading.runs:
        _set_docx_font_color(run, primary)
    if sources:
        for source in sources:
            doc.add_paragraph(source, style="List Bullet")
    else:
        doc.add_paragraph(
            "No se detectaron URLs verificables en el contenido entregado. "
            "Si se usó información externa, el productor debe volver a buscar y citar fuentes."
        )


def _rows_from_json_or_markdown(datos: str) -> tuple[list[str], list[list[Any]], list[str]]:
    sources: list[str] = []
    try:
        parsed = json.loads(datos)
        rows = parsed if isinstance(parsed, list) else [parsed]
        if rows and isinstance(rows[0], dict):
            headers = list(rows[0].keys())
            return headers, [[row.get(header) for header in headers] for row in rows], sources
        normalized = [row if isinstance(row, list) else [row] for row in rows]
        width = max((len(row) for row in normalized), default=1)
        return [f"Columna {index}" for index in range(1, width + 1)], normalized, sources
    except Exception:
        blocks = parse_markdown_blocks(datos)
        sources = aplicar_diseno_entregable(datos, "xlsx")["sources"]
        for block in blocks:
            if block["type"] == "table":
                return block["headers"], block["rows"], sources
        extracted = [[block.get("type", ""), block.get("text", "")] for block in blocks if block.get("text")]
        return ["Tipo", "Contenido"], extracted or [["texto", clean_inline_text(datos)]], sources


def leer_archivo_inteligente(path: str, max_chars: int = 8000) -> str:
    """Lee un archivo local según su tipo y devuelve inventario + extracto."""
    try:
        resolved = resolve_safe_path(path)
    except ValueError as exc:
        return f"Error: {exc}"
    if not resolved.exists():
        return f"Error: '{resolved}' no existe."

    suffix = resolved.suffix.lower()
    if suffix == ".docx":
        return analizar_docx(str(resolved), max_chars=max_chars)
    if suffix == ".pdf":
        return analizar_pdf(str(resolved), max_chars=max_chars)
    if suffix in {".xlsx", ".xlsm", ".csv", ".tsv"}:
        return analizar_xlsx(str(resolved), max_rows=20)
    if suffix == ".pptx":
        return analizar_pptx(str(resolved), max_chars=max_chars)
    return f"Archivo: {resolved}\nTipo: {suffix or 'texto'}\n\n{_read_text_like(resolved, max_chars)}"


def analizar_docx(path: str, max_chars: int = 12000) -> str:
    """Extrae texto y tablas básicas de un documento DOCX."""
    try:
        resolved = resolve_safe_path(path)
        from docx import Document

        doc = Document(str(resolved))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table in doc.tables[:5]:
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows[:10]]
            tables.extend(rows)
        text = "\n".join(paragraphs + tables)
        return f"DOCX: {resolved}\nPárrafos: {len(paragraphs)}\nTablas: {len(doc.tables)}\n\n{text[:max_chars]}"
    except Exception as exc:
        return f"Error al analizar DOCX: {exc}"


def crear_docx(contenido_markdown: str, nombre_sugerido: str = "documento", design_md_path: str = "") -> str:
    """Crea un DOCX nuevo en exports con estilos nativos y diseño DESIGN.md."""
    try:
        from docx import Document
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.shared import Inches

        deliverable = aplicar_diseno_entregable(contenido_markdown, "docx", design_md_path)
        blocks = deliverable["blocks"]
        primary = _theme_color(deliverable, "colors.primary", "#1f5f8b")
        header_text = _theme_color(deliverable, "colors.surface", "#ffffff")
        title = _first_heading(blocks, nombre_sugerido)
        subtitle = next((block["text"] for block in blocks if block["type"] == "paragraph"), "")
        doc = Document()
        doc.core_properties.title = title
        doc.core_properties.subject = "Entregable ejecutivo K-Zero"
        doc.core_properties.author = "K-Zero"
        doc.core_properties.keywords = "K-Zero, DESIGN.md, entregable, ejecutivo"
        section = doc.sections[0]
        section.top_margin = Inches(0.7)
        section.bottom_margin = Inches(0.7)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)
        _configure_docx_styles(doc, primary)
        _add_docx_cover(doc, title, subtitle, primary)

        skipped_cover_title = False
        for block in blocks:
            kind = block["type"]
            if kind == "heading":
                level = int(block.get("level", 1))
                if level == 1 and block["text"] == title and not skipped_cover_title:
                    skipped_cover_title = True
                    continue
                paragraph = doc.add_heading(block["text"], level=0 if level == 1 else min(level - 1, 4))
                for run in paragraph.runs:
                    _set_docx_font_color(run, primary)
            elif kind == "paragraph":
                doc.add_paragraph(block["text"])
            elif kind == "bullets":
                for item in block["items"]:
                    doc.add_paragraph(item, style="List Bullet")
            elif kind == "numbers":
                for item in block["items"]:
                    doc.add_paragraph(item, style="List Number")
            elif kind == "table":
                rows = _table_rows(block)
                table = doc.add_table(rows=len(rows), cols=len(block["headers"]))
                table.style = "Table Grid"
                table.autofit = True
                table.alignment = WD_TABLE_ALIGNMENT.CENTER
                for row_index, row in enumerate(rows):
                    for col_index, value in enumerate(row):
                        cell = table.cell(row_index, col_index)
                        cell.text = str(value)
                        if row_index == 0:
                            _set_docx_cell_shading(cell, primary)
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.bold = True
                                    _set_docx_font_color(run, header_text)
                doc.add_paragraph("")
            elif kind == "code":
                doc.add_paragraph(block["text"], style="Intense Quote")
        _add_docx_sources(doc, deliverable["sources"], primary)
        output = _export_path(nombre_sugerido, ".docx")
        doc.save(str(output))
        return _created_result(output, "docx", nombre_sugerido)
    except Exception as exc:
        return f"Error al crear DOCX: {exc}"


def editar_docx_copia(path: str, instrucciones: str) -> str:
    """Crea una copia DOCX con una nota de edición; nunca modifica el original."""
    try:
        resolved = resolve_safe_path(path)
        from docx import Document

        doc = Document(str(resolved))
        doc.add_heading("Notas de edición solicitadas", level=1)
        doc.add_paragraph(instrucciones)
        output = _export_path(resolved.stem + "_editado", ".docx")
        doc.save(str(output))
        return f"Copia editada: {output}"
    except Exception as exc:
        return f"Error al editar DOCX: {exc}"


def analizar_pdf(path: str, max_chars: int = 12000) -> str:
    """Extrae texto básico de un PDF."""
    try:
        resolved = resolve_safe_path(path)
        text = extract_text(str(resolved))
        return f"PDF: {resolved}\nCaracteres extraídos: {len(text)}\n\n{text[:max_chars]}"
    except Exception as exc:
        return f"Error al analizar PDF: {exc}"


def crear_pdf(contenido_markdown: str, nombre_sugerido: str = "documento", design_md_path: str = "") -> str:
    """Crea un PDF en exports con párrafos, tablas y saltos nativos."""
    try:
        from xml.sax.saxutils import escape

        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle

        output = _export_path(nombre_sugerido, ".pdf")
        deliverable = aplicar_diseno_entregable(contenido_markdown, "pdf", design_md_path)
        blocks = deliverable["blocks"]
        primary = _theme_color(deliverable, "colors.primary", "#1f5f8b")
        border = _theme_color(deliverable, "colors.border", "#d9e2ec")
        styles = getSampleStyleSheet()
        body = ParagraphStyle(
            "KZeroBody",
            parent=styles["BodyText"],
            fontName="Helvetica",
            fontSize=10,
            leading=14,
            spaceAfter=8,
        )
        h1 = ParagraphStyle("KZeroH1", parent=styles["Title"], textColor=colors.HexColor(primary), fontSize=22, leading=27)
        h2 = ParagraphStyle("KZeroH2", parent=styles["Heading2"], textColor=colors.HexColor(primary), fontSize=14, leading=18)
        bullet = ParagraphStyle("KZeroBullet", parent=body, leftIndent=14, firstLineIndent=-8)
        story: list[Any] = []
        for block in blocks:
            kind = block["type"]
            if kind == "heading":
                story.append(Paragraph(escape(block["text"]), h1 if block.get("level") == 1 else h2))
                story.append(Spacer(1, 6))
            elif kind == "paragraph":
                story.append(Paragraph(escape(block["text"]), body))
            elif kind == "bullets":
                for item in block["items"]:
                    story.append(Paragraph(f"• {escape(item)}", bullet))
            elif kind == "numbers":
                for index, item in enumerate(block["items"], start=1):
                    story.append(Paragraph(f"{index}. {escape(item)}", bullet))
            elif kind == "table":
                rows = [
                    [Paragraph(escape(str(cell)), body) for cell in row]
                    for row in _table_rows(block)
                ]
                table = Table(rows, repeatRows=1, hAlign="LEFT")
                table.setStyle(
                    TableStyle(
                        [
                            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(primary)),
                            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(border)),
                            ("VALIGN", (0, 0), (-1, -1), "TOP"),
                            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
                            ("LEFTPADDING", (0, 0), (-1, -1), 6),
                            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                        ]
                    )
                )
                story.append(table)
                story.append(Spacer(1, 10))
            elif kind == "code":
                story.append(Paragraph(escape(block["text"]), styles["Code"]))
        doc = SimpleDocTemplate(
            str(output),
            pagesize=letter,
            rightMargin=0.65 * inch,
            leftMargin=0.65 * inch,
            topMargin=0.65 * inch,
            bottomMargin=0.65 * inch,
            title=_first_heading(blocks, nombre_sugerido),
        )
        doc.build(story)
        return _created_result(output, "pdf", nombre_sugerido)
    except Exception as exc:
        return f"Error al crear PDF: {exc}"


def editar_pdf_copia(path: str, instrucciones: str) -> str:
    """Crea una copia PDF y un archivo de instrucciones asociado en exports."""
    try:
        resolved = resolve_safe_path(path)
        output = _export_path(resolved.stem + "_editado", ".pdf")
        shutil.copy2(resolved, output)
        output.with_suffix(".notes.txt").write_text(instrucciones, encoding="utf-8")
        return f"Copia editada: {output}"
    except Exception as exc:
        return f"Error al editar PDF: {exc}"


def dividir_pdf_copia(path: str, paginas: str) -> str:
    """Crea un PDF con las páginas indicadas, sin modificar el original."""
    try:
        resolved = resolve_safe_path(path)
        from pypdf import PdfReader, PdfWriter

        selected = [int(part.strip()) - 1 for part in paginas.split(",") if part.strip().isdigit()]
        reader = PdfReader(str(resolved))
        writer = PdfWriter()
        for index in selected:
            if 0 <= index < len(reader.pages):
                writer.add_page(reader.pages[index])
        output = _export_path(resolved.stem + "_paginas", ".pdf")
        with output.open("wb") as fh:
            writer.write(fh)
        return _created_result(output, "xlsx", nombre_sugerido)
    except Exception as exc:
        return f"Error al dividir PDF: {exc}"


def combinar_pdf_copia(paths: str, nombre_sugerido: str = "pdf_combinado") -> str:
    """Combina PDFs en una copia nueva en exports."""
    try:
        from pypdf import PdfReader, PdfWriter

        writer = PdfWriter()
        for raw_path in paths.split(";"):
            resolved = resolve_safe_path(raw_path.strip())
            reader = PdfReader(str(resolved))
            for page in reader.pages:
                writer.add_page(page)
        output = _export_path(nombre_sugerido, ".pdf")
        with output.open("wb") as fh:
            writer.write(fh)
        return _created_result(output, "xlsx", nombre_sugerido)
    except Exception as exc:
        return f"Error al combinar PDFs: {exc}"


def analizar_xlsx(path: str, max_rows: int = 20) -> str:
    """Resume hojas y primeras filas de XLSX/CSV/TSV."""
    try:
        resolved = resolve_safe_path(path)
        if resolved.suffix.lower() in {".csv", ".tsv"}:
            delimiter = "\t" if resolved.suffix.lower() == ".tsv" else ","
            with resolved.open("r", encoding="utf-8", errors="replace", newline="") as fh:
                rows = list(csv.reader(fh, delimiter=delimiter))[:max_rows]
            return f"Tabla: {resolved}\nFilas mostradas: {len(rows)}\n" + "\n".join(" | ".join(row) for row in rows)

        from openpyxl import load_workbook

        wb = load_workbook(str(resolved), read_only=True, data_only=False)
        lines = [f"XLSX: {resolved}", "Hojas: " + ", ".join(wb.sheetnames)]
        for sheet_name in wb.sheetnames[:5]:
            ws = wb[sheet_name]
            lines.append(f"\n[{sheet_name}]")
            for row in ws.iter_rows(max_row=max_rows, values_only=True):
                lines.append(" | ".join("" if value is None else str(value) for value in row))
        wb.close()
        return "\n".join(lines)
    except Exception as exc:
        return f"Error al analizar XLSX/CSV: {exc}"


def crear_xlsx(datos: str, nombre_sugerido: str = "datos", design_md_path: str = "") -> str:
    """Crea un XLSX estilizado desde JSON, lista de listas o tabla Markdown."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.table import Table, TableStyleInfo

        deliverable = aplicar_diseno_entregable(datos, "xlsx", design_md_path)
        primary = _theme_color(deliverable, "colors.primary", "#1f5f8b").strip("#")
        headers, rows, sources = _rows_from_json_or_markdown(datos)
        wb = Workbook()
        ws = wb.active
        ws.title = "Matriz"
        ws.append(headers)
        for row in rows:
            ws.append(row)

        header_fill = PatternFill("solid", fgColor=primary)
        for cell in ws[1]:
            cell.fill = header_fill
            cell.font = Font(color="FFFFFF", bold=True)
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        for row in ws.iter_rows(min_row=2):
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
        ws.freeze_panes = "A2"
        max_column = max(ws.max_column, 1)
        max_row = max(ws.max_row, 1)
        ws.auto_filter.ref = f"A1:{get_column_letter(max_column)}{max_row}"
        if max_row >= 2:
            table = Table(displayName="TablaMatriz", ref=ws.auto_filter.ref)
            table.tableStyleInfo = TableStyleInfo(
                name="TableStyleMedium2",
                showFirstColumn=False,
                showLastColumn=False,
                showRowStripes=True,
                showColumnStripes=False,
            )
            ws.add_table(table)
        for column in ws.columns:
            width = min(max(len(str(cell.value or "")) for cell in column) + 2, 45)
            ws.column_dimensions[get_column_letter(column[0].column)].width = width

        if sources:
            source_ws = wb.create_sheet("Fuentes")
            source_ws.append(["URL"])
            for url in sources:
                source_ws.append([url])
            source_ws.freeze_panes = "A2"
            source_ws["A1"].fill = header_fill
            source_ws["A1"].font = Font(color="FFFFFF", bold=True)
            source_ws.column_dimensions["A"].width = min(max(len(url) for url in sources) + 2, 80)

        output = _export_path(nombre_sugerido, ".xlsx")
        wb.save(str(output))
        return _created_result(output, "pptx", nombre_sugerido)
    except Exception as exc:
        return f"Error al crear XLSX: {exc}"


def _crear_xlsx_legacy(datos: str, nombre_sugerido: str = "datos") -> str:
    """Implementación previa conservada solo como referencia interna."""
    try:
        from openpyxl import Workbook

        parsed = json.loads(datos)
        rows = parsed if isinstance(parsed, list) else [parsed]
        wb = Workbook()
        ws = wb.active
        if rows and isinstance(rows[0], dict):
            headers = list(rows[0].keys())
            ws.append(headers)
            for row in rows:
                ws.append([row.get(header) for header in headers])
        else:
            for row in rows:
                ws.append(row if isinstance(row, list) else [row])
        output = _export_path(nombre_sugerido, ".xlsx")
        wb.save(str(output))
        return f"Archivo creado: {output}"
    except Exception as exc:
        return f"Error al crear XLSX: {exc}"


def editar_xlsx_copia(path: str, instrucciones: str) -> str:
    """Crea una copia XLSX con hoja de notas de edición."""
    try:
        resolved = resolve_safe_path(path)
        from openpyxl import load_workbook

        wb = load_workbook(str(resolved))
        ws = wb.create_sheet("Notas edición")
        ws.append(["Instrucciones"])
        ws.append([instrucciones])
        output = _export_path(resolved.stem + "_editado", ".xlsx")
        wb.save(str(output))
        return f"Copia editada: {output}"
    except Exception as exc:
        return f"Error al editar XLSX: {exc}"


def analizar_pptx(path: str, max_chars: int = 12000) -> str:
    """Extrae texto básico de una presentación PPTX."""
    try:
        resolved = resolve_safe_path(path)
        from pptx import Presentation

        prs = Presentation(str(resolved))
        lines = [f"PPTX: {resolved}", f"Diapositivas: {len(prs.slides)}"]
        for index, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            lines.append(f"\nSlide {index}: " + " | ".join(texts))
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        return f"Error al analizar PPTX: {exc}"


def crear_pptx(contenido_markdown: str, nombre_sugerido: str = "presentacion", design_md_path: str = "") -> str:
    """Crea una presentación PPTX con slides nativos y diseño consistente."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        deliverable = aplicar_diseno_entregable(contenido_markdown, "pptx", design_md_path)
        blocks = deliverable["blocks"]
        primary = _theme_color(deliverable, "colors.primary", "#1f5f8b")
        surface = _theme_color(deliverable, "colors.surface", "#ffffff")
        red, green, blue = _hex_rgb(primary)
        bg_red, bg_green, bg_blue = _hex_rgb(surface)
        prs = Presentation()
        title = _first_heading(blocks, "Presentación")
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = RGBColor(bg_red, bg_green, bg_blue)
        title_slide.shapes.title.text = title
        title_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(red, green, blue)
        title_slide.placeholders[1].text = "Entregable generado por K-Zero"

        sections: list[dict[str, Any]] = []
        current: dict[str, Any] | None = None
        for block in blocks:
            if block["type"] == "heading" and block.get("level", 1) <= 2:
                if current is not None:
                    sections.append(current)
                current = {"title": block["text"], "blocks": []}
            elif current is not None:
                current["blocks"].append(block)
        if current is not None:
            sections.append(current)
        if not sections:
            sections = [{"title": title, "blocks": blocks}]

        for section_data in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(bg_red, bg_green, bg_blue)
            title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.7), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = section_data["title"]
            title_run = title_frame.paragraphs[0].runs[0]
            title_run.font.bold = True
            title_run.font.size = Pt(24)
            title_run.font.color.rgb = RGBColor(red, green, blue)
            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(8.2), Inches(5.0))
            text_frame = body_box.text_frame
            text_frame.word_wrap = True
            first_paragraph = True
            for block in section_data["blocks"]:
                if block["type"] == "paragraph":
                    paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
                    paragraph.text = block["text"][:450]
                    paragraph.font.size = Pt(16)
                    first_paragraph = False
                elif block["type"] in {"bullets", "numbers"}:
                    for item in block["items"][:7]:
                        paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
                        paragraph.text = item[:180]
                        paragraph.level = 0
                        paragraph.font.size = Pt(15)
                        first_paragraph = False
                elif block["type"] == "table":
                    rows = _table_rows(block)[:6]
                    cols = min(len(block["headers"]), 4)
                    table_shape = slide.shapes.add_table(
                        len(rows),
                        cols,
                        Inches(0.7),
                        Inches(1.35),
                        Inches(8.2),
                        Inches(0.35 * len(rows) + 0.4),
                    )
                    table = table_shape.table
                    for row_index, row in enumerate(rows):
                        for col_index in range(cols):
                            cell = table.cell(row_index, col_index)
                            cell.text = str(row[col_index] if col_index < len(row) else "")[:80]
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10)
                                    if row_index == 0:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                            if row_index == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(red, green, blue)
                    break
            if first_paragraph:
                text_frame.text = "Resumen visual de la sección."
        output = _export_path(nombre_sugerido, ".pptx")
        prs.save(str(output))
        return f"Archivo creado: {output}"
    except Exception as exc:
        return f"Error al crear PPTX: {exc}"


def editar_pptx_copia(path: str, instrucciones: str) -> str:
    """Crea una copia PPTX con una diapositiva de notas de edición."""
    try:
        resolved = resolve_safe_path(path)
        from pptx import Presentation

        prs = Presentation(str(resolved))
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Notas de edición solicitadas"
        slide.placeholders[1].text = instrucciones
        output = _export_path(resolved.stem + "_editado", ".pptx")
        prs.save(str(output))
        return f"Copia editada: {output}"
    except Exception as exc:
        return f"Error al editar PPTX: {exc}"


def analizar_archivos_frontend(path: str, max_chars: int = 12000) -> str:
    """Revisa archivos HTML/CSS/JS locales y resume hallazgos de interfaz."""
    try:
        resolved = resolve_safe_path(path)
        files = [resolved] if resolved.is_file() else [
            file for file in resolved.rglob("*") if file.suffix.lower() in {".html", ".css", ".js", ".jsx", ".ts", ".tsx"}
        ][:30]
        lines = [f"Archivos frontend analizados: {len(files)}"]
        for file in files:
            text = file.read_text(encoding="utf-8", errors="replace")
            lines.append(f"\nArchivo: {file}")
            lines.append(f"- Tamaño: {len(text)} caracteres")
            lines.append(f"- Usa forms: {'<form' in text.lower()}")
            lines.append(f"- Tiene alt=: {'alt=' in text.lower()}")
            lines.append(text[:800])
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        return f"Error al analizar frontend: {exc}"


def validar_entregable(texto: str, requisitos: str) -> str:
    """Valida de forma simple si un texto cumple requisitos explícitos."""
    missing = []
    warnings = []
    if "**" in texto or "###" in texto or any(line.strip().startswith("|") for line in texto.splitlines()):
        warnings.append("El entregable parece contener Markdown crudo visible; usa limpiar_markdown_entregable o render nativo.")
    for raw in requisitos.replace(";", "\n").splitlines():
        requirement = raw.strip("- ").strip()
        if requirement and requirement.lower() not in texto.lower():
            missing.append(requirement)
    if not missing and not warnings:
        return "Validación: cumple los requisitos explícitos revisados."
    lines = ["Validación:"]
    if missing:
        lines.append("Faltan posibles requisitos:")
        lines.extend(f"- {item}" for item in missing[:20])
    if warnings:
        lines.append("Advertencias de formato:")
        lines.extend(f"- {item}" for item in warnings)
    return "\n".join(lines)
