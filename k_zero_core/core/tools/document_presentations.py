"""Tools PPTX para lectura, creación y edición por copia."""
from __future__ import annotations

from typing import Any

from k_zero_core.core.tool_safety import resolve_safe_path
from k_zero_core.core.tools.document_common import (
    _export_path,
    _first_heading,
    _hex_rgb,
    _table_rows,
    _theme_color,
)
from k_zero_core.services.design_md import aplicar_diseno_entregable


def analizar_pptx(path: str, max_chars: int = 12000) -> str:
    """Extrae texto básico de una presentación PPTX."""
    try:
        resolved = resolve_safe_path(path)
        from pptx import Presentation

        prs = Presentation(str(resolved))
        lines = [f"PPTX: {resolved}", f"Diapositivas: {len(prs.slides)}"]
        for index, slide in enumerate(prs.slides, start=1):
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            lines.append(f"\nSlide {index}: " + " | ".join(texts))
        return "\n".join(lines)[:max_chars]
    except Exception as exc:
        return f"Error al analizar PPTX: {exc}"


def crear_pptx(contenido_markdown: str, nombre_sugerido: str = "presentacion", design_md_path: str = "") -> str:
    """Crea una presentación PPTX con slides nativos y diseño consistente."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        deliverable = aplicar_diseno_entregable(contenido_markdown, "pptx", design_md_path)
        blocks = deliverable["blocks"]
        primary = _theme_color(deliverable, "colors.primary", "#1f5f8b")
        surface = _theme_color(deliverable, "colors.surface", "#ffffff")
        red, green, blue = _hex_rgb(primary)
        bg_red, bg_green, bg_blue = _hex_rgb(surface)
        prs = Presentation()
        title = _first_heading(blocks, "Presentación")
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = RGBColor(bg_red, bg_green, bg_blue)
        title_slide.shapes.title.text = title
        title_slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(red, green, blue)
        title_slide.placeholders[1].text = "Entregable generado por K-Zero"

        sections: list[dict[str, Any]] = []
        current: dict[str, Any] | None = None
        for block in blocks:
            if block["type"] == "heading" and block.get("level", 1) <= 2:
                if current is not None:
                    sections.append(current)
                current = {"title": block["text"], "blocks": []}
            elif current is not None:
                current["blocks"].append(block)
        if current is not None:
            sections.append(current)
        if not sections:
            sections = [{"title": title, "blocks": blocks}]

        for section_data in sections:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(bg_red, bg_green, bg_blue)
            title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(8.7), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = section_data["title"]
            title_run = title_frame.paragraphs[0].runs[0]
            title_run.font.bold = True
            title_run.font.size = Pt(24)
            title_run.font.color.rgb = RGBColor(red, green, blue)
            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.25), Inches(8.2), Inches(5.0))
            text_frame = body_box.text_frame
            text_frame.word_wrap = True
            first_paragraph = True
            for block in section_data["blocks"]:
                if block["type"] == "paragraph":
                    paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
                    paragraph.text = block["text"][:450]
                    paragraph.font.size = Pt(16)
                    first_paragraph = False
                elif block["type"] in {"bullets", "numbers"}:
                    for item in block["items"][:7]:
                        paragraph = text_frame.paragraphs[0] if first_paragraph else text_frame.add_paragraph()
                        paragraph.text = item[:180]
                        paragraph.level = 0
                        paragraph.font.size = Pt(15)
                        first_paragraph = False
                elif block["type"] == "table":
                    rows = _table_rows(block)[:6]
                    cols = min(len(block["headers"]), 4)
                    table_shape = slide.shapes.add_table(
                        len(rows),
                        cols,
                        Inches(0.7),
                        Inches(1.35),
                        Inches(8.2),
                        Inches(0.35 * len(rows) + 0.4),
                    )
                    table = table_shape.table
                    for row_index, row in enumerate(rows):
                        for col_index in range(cols):
                            cell = table.cell(row_index, col_index)
                            cell.text = str(row[col_index] if col_index < len(row) else "")[:80]
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10)
                                    if row_index == 0:
                                        run.font.bold = True
                                        run.font.color.rgb = RGBColor(255, 255, 255)
                            if row_index == 0:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(red, green, blue)
                    break
            if first_paragraph:
                text_frame.text = "Resumen visual de la sección."
        output = _export_path(nombre_sugerido, ".pptx")
        prs.save(str(output))
        return f"Archivo creado: {output}"
    except Exception as exc:
        return f"Error al crear PPTX: {exc}"


def editar_pptx_copia(path: str, instrucciones: str) -> str:
    """Crea una copia PPTX con una diapositiva de notas de edición."""
    try:
        resolved = resolve_safe_path(path)
        from pptx import Presentation

        prs = Presentation(str(resolved))
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Notas de edición solicitadas"
        slide.placeholders[1].text = instrucciones
        output = _export_path(resolved.stem + "_editado", ".pptx")
        prs.save(str(output))
        return f"Copia editada: {output}"
    except Exception as exc:
        return f"Error al editar PPTX: {exc}"
